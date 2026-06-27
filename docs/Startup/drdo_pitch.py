"""
InfraCore Systems — DRDO-ADE Pitch Deck (4 slides)
Generates: docs/Startup/InfraCore_DRDO_Pitch.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Paths ────────────────────────────────────────────────────────────────────
DIR   = os.path.dirname(os.path.abspath(__file__))
OUT   = os.path.join(DIR, 'InfraCore_DRDO_Pitch.pptx')

# ── Colours ──────────────────────────────────────────────────────────────────
def col(h): return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

NAVY   = col('1B3A6B')
BLUE   = col('2E75B6')
TEAL   = col('0D7377')
WHITE  = col('FFFFFF')
LGRAY  = col('F2F4F8')
DGRAY  = col('1A1A1A')
SILVER = col('D0D7E4')
AMBER  = col('E67E22')
GREEN  = col('27AE60')

# ── Slide dimensions (16:9) ──────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

# ── Helpers ──────────────────────────────────────────────────────────────────
def add_slide(prs):
    blank = prs.slide_layouts[6]          # fully blank
    slide = prs.slides.add_slide(blank)
    slide.shapes.title                     # may not exist on blank — ignore
    return slide

def box(slide, x, y, w, h, bg=None, border=None, border_w=1):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.width = Pt(border_w)
    if bg:
        shape.fill.solid(); shape.fill.fore_color.rgb = bg
    else:
        shape.fill.background()
    if border:
        shape.line.color.rgb = border
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, x, y, w, h,
        size=18, bold=False, color=DGRAY, align=PP_ALIGN.LEFT,
        wrap=True, italic=False):
    tf = slide.shapes.add_textbox(x, y, w, h)
    tf.word_wrap = wrap
    p  = tf.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return tf

def bullet_block(slide, items, x, y, w, h,
                 size=16, color=DGRAY, gap=Pt(10), bold_first=False):
    """items: list of str or (heading, body) tuples"""
    tf = slide.shapes.add_textbox(x, y, w, h)
    tf.word_wrap = True
    tf.text_frame.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.text_frame.paragraphs[0]
            first = False
        else:
            p = tf.text_frame.add_paragraph()
        p.space_before = gap
        if isinstance(item, tuple):
            head, body = item
            run = p.add_run()
            run.text = f'•  {head}'
            run.font.size  = Pt(size)
            run.font.bold  = True
            run.font.color.rgb = color
            p2 = tf.text_frame.add_paragraph()
            p2.space_before = Pt(2)
            r2 = p2.add_run()
            r2.text = f'    {body}'
            r2.font.size = Pt(size - 1)
            r2.font.color.rgb = color
        else:
            run = p.add_run()
            run.text = f'•  {item}'
            run.font.size = Pt(size)
            run.font.bold = bold_first and first
            run.font.color.rgb = color
    return tf

def top_bar(slide, title, subtitle=None):
    """Navy bar across top with white title."""
    box(slide, 0, 0, W, Inches(1.15), bg=NAVY)
    txt(slide, title,
        Inches(0.35), Inches(0.12), Inches(12), Inches(0.72),
        size=30, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        txt(slide, subtitle,
            Inches(0.35), Inches(0.78), Inches(12), Inches(0.35),
            size=14, bold=False, color=SILVER, align=PP_ALIGN.LEFT)

def divider(slide, y, color=BLUE, thickness=2):
    line = slide.shapes.add_shape(1, Inches(0.35), y, Inches(12.6), Pt(thickness))
    line.fill.solid(); line.fill.fore_color.rgb = color
    line.line.fill.background()

def card(slide, x, y, w, h, title, lines, title_bg=NAVY, body_bg=LGRAY,
         tsize=14, bsize=13):
    # Header
    box(slide, x, y, w, Inches(0.42), bg=title_bg)
    txt(slide, title, x + Inches(0.12), y + Inches(0.05),
        w - Inches(0.25), Inches(0.35),
        size=tsize, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    # Body
    bh = h - Inches(0.42)
    box(slide, x, y + Inches(0.42), w, bh, bg=body_bg, border=SILVER, border_w=1)
    content = '\n'.join(f'•  {l}' for l in lines)
    txt(slide, content,
        x + Inches(0.15), y + Inches(0.50),
        w - Inches(0.3), bh - Inches(0.15),
        size=bsize, color=DGRAY, wrap=True)

def pic(slide, fname, x, y, w, h):
    path = os.path.join(DIR, fname)
    if os.path.exists(path):
        slide.shapes.add_picture(path, x, y, w, h)
    else:
        b = box(slide, x, y, w, h, bg=LGRAY, border=SILVER)
        txt(slide, f'[{fname}]', x, y + h/2 - Inches(0.15),
            w, Inches(0.3), size=10, color=SILVER, align=PP_ALIGN.CENTER)

# ============================================================================
#  SLIDE 1  —  TITLE
# ============================================================================
def slide_title(prs):
    slide = add_slide(prs)

    # Full navy background
    box(slide, 0, 0, W, H, bg=NAVY)

    # Accent bar top
    box(slide, 0, 0, W, Inches(0.08), bg=BLUE)

    # Company name
    txt(slide, 'INFRACORE SYSTEMS LLP',
        Inches(0.7), Inches(1.1), Inches(11), Inches(0.55),
        size=16, bold=True, color=SILVER, align=PP_ALIGN.LEFT)

    # Main title
    txt(slide, 'Precision Guidance\nfor GPS-Denied Environments',
        Inches(0.7), Inches(1.7), Inches(9.5), Inches(2.0),
        size=40, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Subtitle
    txt(slide, 'An indigenous, FPGA-accelerated terminal-phase tracking system\nbuilt on the SHAKTI RISC-V platform',
        Inches(0.7), Inches(3.75), Inches(9.5), Inches(0.85),
        size=16, bold=False, color=SILVER, align=PP_ALIGN.LEFT)

    # Divider line
    divider(slide, Inches(4.72), color=BLUE, thickness=2)

    # Team details bottom-left
    details = [
        ('Team',        'Risky 5  —  InfraCore Systems LLP'),
        ('Institution', 'Shiv Nadar Institution of Eminence'),
        ('Mentor',      'Dr. Venkatnarayan Hariharan'),
    ]
    y0 = Inches(4.9)
    for label, val in details:
        txt(slide, label.upper(),
            Inches(0.7), y0, Inches(1.5), Inches(0.35),
            size=9, bold=True, color=col('5B8DB8'), align=PP_ALIGN.LEFT)
        txt(slide, val,
            Inches(2.2), y0, Inches(7.5), Inches(0.35),
            size=12, bold=False, color=WHITE, align=PP_ALIGN.LEFT)
        y0 += Inches(0.42)

    # Logo placeholder bottom-right
    txt(slide, 'DIR-V  |  April 2026',
        Inches(9.5), Inches(6.9), Inches(3.5), Inches(0.4),
        size=11, bold=False, color=col('5B8DB8'), align=PP_ALIGN.RIGHT)

# ============================================================================
#  SLIDE 2  —  THE PROBLEM & WHAT WE'RE DOING DIFFERENTLY
# ============================================================================
def slide_problem(prs):
    slide = add_slide(prs)
    top_bar(slide, 'The Problem', 'Why terminal-phase visual guidance matters')

    # ── Left column: problem ─────────────────────────────────────────────────
    txt(slide, 'THE GAP',
        Inches(0.35), Inches(1.35), Inches(6), Inches(0.35),
        size=11, bold=True, color=BLUE)

    problems = [
        ('GPS is not reliable in the terminal phase',
         'Adversarial jammers are deployed precisely at the last moment. A missile relying only on GPS coordinates becomes blind exactly when it matters most.'),
        ('Existing seekers rely on foreign compute',
         'SCALP-EG, HAMMER, and similar IIR seekers run on proprietary foreign black-box processors — creating supply-chain risk and ITAR dependency.'),
        ('Standard IR sensors only track hot targets',
         'Conventional IR seekers need an active heat source (jet exhaust, engine). Stationary infrastructure, grounded vehicles, and inanimate targets are invisible to them.'),
    ]

    y = Inches(1.75)
    for head, body in problems:
        # accent dot
        box(slide, Inches(0.35), y + Inches(0.08), Inches(0.06), Inches(0.06), bg=BLUE)
        txt(slide, head,
            Inches(0.55), y, Inches(5.9), Inches(0.35),
            size=14, bold=True, color=NAVY)
        txt(slide, body,
            Inches(0.55), y + Inches(0.35), Inches(5.85), Inches(0.72),
            size=12, color=DGRAY)
        y += Inches(1.18)

    # ── Vertical divider ─────────────────────────────────────────────────────
    div = slide.shapes.add_shape(1, Inches(6.6), Inches(1.3), Pt(2), Inches(5.9))
    div.fill.solid(); div.fill.fore_color.rgb = SILVER
    div.line.fill.background()

    # ── Right column: our answer ─────────────────────────────────────────────
    txt(slide, 'OUR ANSWER',
        Inches(6.85), Inches(1.35), Inches(6), Inches(0.35),
        size=11, bold=True, color=TEAL)

    answers = [
        ('100% Indigenous Hardware',
         'SHAKTI Yamuna RISC-V processor + Artix-7 FPGA — no foreign CPUs, no black boxes, no ITAR constraints.'),
        ('Works when GPS is Jammed',
         'Onboard image tracking takes over autonomously in the terminal phase. The system locks on visually and guides itself to target.'),
        ('LWIR: Tracks Thermally Passive Targets',
         'Long-Wave IR detects intrinsic heat of any object — structures, vehicles at rest, infrastructure. No active heat source required.'),
        ('Real-Time at < 200 microseconds',
         'FPGA-accelerated correlation: 19,300 clock cycles per tracking decision at 100 MHz. Faster than any software solution.'),
    ]

    y = Inches(1.75)
    for head, body in answers:
        box(slide, Inches(6.85), y + Inches(0.08), Inches(0.06), Inches(0.06), bg=TEAL)
        txt(slide, head,
            Inches(7.05), y, Inches(5.9), Inches(0.35),
            size=14, bold=True, color=NAVY)
        txt(slide, body,
            Inches(7.05), y + Inches(0.35), Inches(5.85), Inches(0.62),
            size=12, color=DGRAY)
        y += Inches(1.05)

    # Bottom bar
    box(slide, 0, Inches(7.1), W, Inches(0.4), bg=LGRAY)
    txt(slide, 'Modular & sensor-agnostic — same platform adapts to UAVs, loitering munitions, border surveillance',
        Inches(0.35), Inches(7.13), Inches(12.6), Inches(0.32),
        size=11, bold=False, color=NAVY, align=PP_ALIGN.CENTER, italic=True)

# ============================================================================
#  SLIDE 3  —  CORE TECHNOLOGY
# ============================================================================
def slide_tech(prs):
    slide = add_slide(prs)
    top_bar(slide, 'Core Technology', 'FPGA-accelerated image correlation on indigenous hardware')

    # ── Architecture image left ───────────────────────────────────────────────
    pic(slide, 'architecture.png', Inches(0.3), Inches(1.25), Inches(5.8), Inches(3.8))

    txt(slide, 'Yamuna RTOS Application — SHAKTI processor manages sensor fusion, queues results, and dispatches guidance commands.',
        Inches(0.3), Inches(5.1), Inches(5.8), Inches(0.6),
        size=10, color=col('555555'), italic=True)

    # ── Right: how it works ───────────────────────────────────────────────────
    txt(slide, 'HOW IT WORKS',
        Inches(6.35), Inches(1.3), Inches(6.6), Inches(0.35),
        size=11, bold=True, color=BLUE)

    steps = [
        ('1.  Acquire',  'Camera/LWIR frame arrives. NCC (Normalised Cross-Correlation) scans the full image to find and lock on to the target.'),
        ('2.  Track',    'KCF algorithm tracks the target frame-to-frame in the frequency domain — no re-scanning needed. Fixed compute cost regardless of image size.'),
        ('3.  Decide',   'FPGA writes displacement (Δx, Δy) to a register. SHAKTI RTOS reads it, fuses with IMU data, and outputs a guidance command.'),
        ('4.  Correct',  'Flight controller receives YAW / PITCH corrections and adjusts heading — closing the loop autonomously.'),
    ]

    y = Inches(1.72)
    for label, body in steps:
        # Step badge
        badge = box(slide, Inches(6.35), y, Inches(1.35), Inches(0.32), bg=NAVY)
        txt(slide, label,
            Inches(6.38), y + Inches(0.02), Inches(1.3), Inches(0.28),
            size=12, bold=True, color=WHITE)
        txt(slide, body,
            Inches(7.78), y, Inches(5.2), Inches(0.65),
            size=12, color=DGRAY)
        y += Inches(0.82)

    # ── Key numbers bottom row ────────────────────────────────────────────────
    divider(slide, Inches(5.15), color=SILVER, thickness=1)

    metrics = [
        ('193 µs',        'Per tracking\ndecision'),
        ('32 × 32',       'Patch size\n(scalable to 64×64)'),
        ('5,000+',        'Decisions\nper second'),
        ('100% Indigenous','SHAKTI + FPGA\nno foreign silicon'),
    ]

    xpos = [Inches(0.3), Inches(3.45), Inches(6.6), Inches(9.75)]
    for (val, label), x in zip(metrics, xpos):
        bx = box(slide, x, Inches(5.3), Inches(3.0), Inches(1.85), bg=LGRAY, border=SILVER)
        txt(slide, val,
            x + Inches(0.12), Inches(5.4), Inches(2.75), Inches(0.75),
            size=26, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        txt(slide, label,
            x + Inches(0.12), Inches(6.1), Inches(2.75), Inches(0.65),
            size=12, color=col('555555'), align=PP_ALIGN.CENTER)

# ============================================================================
#  SLIDE 4  —  WHERE WE ARE & WHAT WE'RE LOOKING FOR
# ============================================================================
def slide_status(prs):
    slide = add_slide(prs)
    top_bar(slide, 'Where We Are', 'Status, demonstrated results, and what we need next')

    # ── Left: what's done ────────────────────────────────────────────────────
    txt(slide, 'DEMONSTRATED',
        Inches(0.35), Inches(1.32), Inches(6), Inches(0.32),
        size=11, bold=True, color=GREEN)

    done_items = [
        'KCF/MOSSE tracking RTL verified in Vivado simulation — outputs match Python reference model exactly',
        'End-to-end frame capture: camera module → SHAKTI Yamuna SoC (serial logs on physical hardware)',
        'NCC → KCF handoff pipeline: target selected on Frame 1, found by NCC on Frame 2, tracked by KCF from Frame 3',
        'AXI4-Lite register interface designed — SHAKTI reads displacement and TARGET_FOUND flag from FPGA',
        'RTOS application architecture: task queues, ISR, control logic, flight command dispatch',
    ]

    y = Inches(1.72)
    for item in done_items:
        box(slide, Inches(0.35), y + Inches(0.1), Inches(0.12), Inches(0.12), bg=GREEN)
        txt(slide, item,
            Inches(0.57), y, Inches(5.9), Inches(0.55),
            size=12, color=DGRAY)
        y += Inches(0.62)

    # Demo image
    pic(slide, 'demo_s6_summary.png',
        Inches(0.35), Inches(5.1), Inches(5.8), Inches(2.05))

    # ── Vertical divider ─────────────────────────────────────────────────────
    div = slide.shapes.add_shape(1, Inches(6.6), Inches(1.3), Pt(2), Inches(5.9))
    div.fill.solid(); div.fill.fore_color.rgb = SILVER
    div.line.fill.background()

    # ── Right: next steps + ask ───────────────────────────────────────────────
    txt(slide, 'NEXT MILESTONES',
        Inches(6.85), Inches(1.32), Inches(6), Inches(0.32),
        size=11, bold=True, color=AMBER)

    next_items = [
        'Continuous video streaming (>10 FPS) through FPGA pipeline',
        'LWIR thermal sensor integration (replacing current camera surrogate)',
        'IMU sensor fusion on SHAKTI RTOS',
        'Full hardware close-loop test on physical FPGA board',
    ]

    y = Inches(1.72)
    for item in next_items:
        box(slide, Inches(6.85), y + Inches(0.1), Inches(0.12), Inches(0.12), bg=AMBER)
        txt(slide, item,
            Inches(7.07), y, Inches(5.85), Inches(0.55),
            size=12, color=DGRAY)
        y += Inches(0.62)

    # Ask box
    box(slide, Inches(6.85), Inches(4.45), Inches(6.15), Inches(2.8), bg=NAVY)

    txt(slide, 'WHAT WE ARE LOOKING FOR',
        Inches(7.05), Inches(4.6), Inches(5.75), Inches(0.38),
        size=12, bold=True, color=SILVER)

    ask_items = [
        'Technical feedback from domain experts on image-based terminal guidance',
        'Access to test environments or UAV/munition simulation platforms',
        'Potential collaboration or field validation opportunities',
        'Guidance on iDEX / TDF funding pathways for next phase',
    ]

    y = Inches(5.05)
    for item in ask_items:
        txt(slide, f'•  {item}',
            Inches(7.05), y, Inches(5.75), Inches(0.48),
            size=12, color=WHITE)
        y += Inches(0.52)

# ============================================================================
#  MAIN
# ============================================================================
def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    print('Building slides...')
    slide_title(prs)   ; print('  Slide 1: Title')
    slide_problem(prs) ; print('  Slide 2: Problem & Differentiation')
    slide_tech(prs)    ; print('  Slide 3: Core Technology')
    slide_status(prs)  ; print('  Slide 4: Status & Ask')

    prs.save(OUT)
    print(f'\nSaved: {OUT}  ({prs.slides.__len__()} slides)')

if __name__ == '__main__':
    main()
