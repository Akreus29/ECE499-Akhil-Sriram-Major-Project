"""
make_presentation.py
Generates docs/KCF_Final_Presentation.pptx
Professional light theme — 16 slides
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
ROOT_DIR   = os.path.join(SCRIPT_DIR, '..')
IMG_DIR    = os.path.join(ROOT_DIR, 'docs', 'MidTerm Report')
OUT_PATH   = os.path.join(ROOT_DIR, 'docs', 'KCF_Final_Presentation.pptx')

# ── Colour palette ─────────────────────────────────────────────────────────────
def col(h): return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

NAVY   = col('1B3A6B')
BLUE   = col('2E75B6')
DTEXT  = col('1A1A1A')
WHITE  = col('FFFFFF')
LGRAY  = col('F2F4F8')
SILVER = col('D0D7E4')

SW, SH = Inches(13.33), Inches(7.5)


# ── Low-level helpers ──────────────────────────────────────────────────────────

def _run(para, text, size=18, bold=False, color=None, italic=False):
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color or DTEXT
    return run

def _fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb

def _border(shape, rgb, width_pt=1.0):
    ln = shape.line
    ln.color.rgb = rgb
    ln.width = Pt(width_pt)

def tb(slide, text, x, y, w, h,
       size=18, bold=False, color=None, align=PP_ALIGN.LEFT,
       fill=None, border=None, italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = align
    _run(para, text, size=size, bold=bold, color=color or DTEXT, italic=italic)
    if fill:  _fill(txb, fill)
    if border: _border(txb, border)
    return txb

def hbar(slide, y, x=None, w=None, color=None, height=Inches(0.04)):
    x = x if x is not None else Inches(0.3)
    w = w or (SW - Inches(0.6))
    bar = slide.shapes.add_shape(1, x, y, w, height)
    _fill(bar, color or NAVY)
    bar.line.fill.background()
    return bar

def pic(slide, fname, x, y, w, h):
    path = os.path.join(IMG_DIR, fname) if not os.path.isabs(fname) else fname
    if not os.path.exists(path):
        p = slide.shapes.add_textbox(x, y, w, h).text_frame.paragraphs[0]
        _run(p, f'[IMAGE: {os.path.basename(path)}]', size=10, color=col('CC0000'))
        return None
    return slide.shapes.add_picture(path, x, y, w, h)

def bullets(slide, items, x, y, w, h, size=16, color=None, bullet='•', leading=4):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        para.space_before = Pt(leading)
        if isinstance(item, tuple):
            head, body = item
            _run(para, f'{bullet} ', size=size, bold=True,  color=NAVY)
            _run(para, head,          size=size, bold=True,  color=NAVY)
            _run(para, f'  {body}',   size=size, bold=False, color=color or DTEXT)
        else:
            _run(para, f'{bullet}  {item}', size=size, color=color or DTEXT)
        first = False
    return txb

def sec(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    bar = slide.shapes.add_shape(1, Inches(2.5), Inches(3.6), Inches(8.33), Inches(0.05))
    _fill(bar, BLUE); bar.line.fill.background()
    tb(slide, title, Inches(0.5), Inches(2.7), Inches(12.33), Inches(1.6),
       size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    return slide

def white_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    return slide

def slide_header(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), SW, Inches(1.1))
    _fill(bar, NAVY); bar.line.fill.background()
    tb(slide, title, Inches(0.35), Inches(0.1), Inches(10.5), Inches(0.9),
       size=26, bold=True, color=WHITE)
    if subtitle:
        tb(slide, subtitle, Inches(10.9), Inches(0.3), Inches(2.1), Inches(0.5),
           size=13, color=col('A0BAD8'), align=PP_ALIGN.RIGHT)

def card(slide, x, y, w, h, title, body_lines, title_size=15, body_size=13):
    hd = slide.shapes.add_shape(1, x, y, w, Inches(0.42))
    _fill(hd, NAVY); hd.line.fill.background()
    tb(slide, title, x + Inches(0.12), y + Inches(0.05),
       w - Inches(0.15), Inches(0.35), size=title_size, bold=True, color=WHITE)
    bd = slide.shapes.add_shape(1, x, y + Inches(0.42), w, h - Inches(0.42))
    _fill(bd, LGRAY); _border(bd, SILVER, 0.5)
    txb = slide.shapes.add_textbox(x + Inches(0.12), y + Inches(0.5),
                                   w - Inches(0.2), h - Inches(0.6))
    tf = txb.text_frame; tf.word_wrap = True
    first = True
    for line in body_lines:
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        para.space_before = Pt(2)
        _run(para, f'•  {line}', size=body_size, color=DTEXT)
        first = False

def tbl(slide, headers, rows, x, y, w, h, hdr_fill=None, font_size=13):
    nc = len(headers); nr = len(rows) + 1
    cw = w / nc; rh = h / nr
    hf = hdr_fill or NAVY
    for ci, hdr in enumerate(headers):
        cx = x + ci * cw
        cell = slide.shapes.add_shape(1, cx, y, cw, rh)
        _fill(cell, hf); _border(cell, WHITE, 0.5)
        tb(slide, hdr, cx + Inches(0.05), y + Inches(0.04),
           cw - Inches(0.1), rh - Inches(0.06),
           size=font_size, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for ri, row in enumerate(rows):
        ry = y + (ri + 1) * rh
        fill = LGRAY if ri % 2 == 0 else WHITE
        for ci, cell_text in enumerate(row):
            cx = x + ci * cw
            cell = slide.shapes.add_shape(1, cx, ry, cw, rh)
            _fill(cell, fill); _border(cell, SILVER, 0.3)
            tb(slide, cell_text, cx + Inches(0.06), ry + Inches(0.04),
               cw - Inches(0.1), rh - Inches(0.06),
               size=font_size - 1, bold=(ci == 0),
               color=NAVY if ci == 0 else DTEXT,
               align=PP_ALIGN.CENTER)


# ── Slides ─────────────────────────────────────────────────────────────────────

def slide_title(prs):
    slide = white_slide(prs)
    for bar_y in [Inches(0), Inches(7.32)]:
        b = slide.shapes.add_shape(1, Inches(0), bar_y, SW, Inches(0.18))
        _fill(b, NAVY); b.line.fill.background()
    acc = slide.shapes.add_shape(1, Inches(0), Inches(0.18), Inches(0.18), Inches(7.14))
    _fill(acc, BLUE); acc.line.fill.background()

    tb(slide, 'RTL Implementation of the',
       Inches(0.5), Inches(1.15), Inches(12.5), Inches(0.8),
       size=34, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    tb(slide, 'KCF / MOSSE Visual Tracking Algorithm on FPGA',
       Inches(0.5), Inches(1.88), Inches(12.5), Inches(0.9),
       size=34, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    hbar(slide, Inches(2.93), x=Inches(1.5), w=Inches(10.33), color=BLUE, height=Inches(0.035))
    tb(slide, 'Akhil Sriram  |  Roll No. 2210110134  |  B.Tech ECE',
       Inches(0.5), Inches(3.1), Inches(12.5), Inches(0.5),
       size=18, color=DTEXT, align=PP_ALIGN.CENTER)
    tb(slide, 'Supervisor: Dr. Venkatnarayan Hariharan',
       Inches(0.5), Inches(3.58), Inches(12.5), Inches(0.45),
       size=16, color=col('555555'), align=PP_ALIGN.CENTER)
    tb(slide, 'Department of Electrical Engineering  |  Shiv Nadar Institution of Eminence',
       Inches(0.5), Inches(3.98), Inches(12.5), Inches(0.45),
       size=14, color=col('666666'), align=PP_ALIGN.CENTER)
    tb(slide, 'April 2026',
       Inches(0.5), Inches(4.38), Inches(12.5), Inches(0.4),
       size=13, italic=True, color=col('888888'), align=PP_ALIGN.CENTER)
    logo = os.path.join(IMG_DIR, 'End_Term_Report_Template', 'logonew.jpg')
    if os.path.exists(logo):
        slide.shapes.add_picture(logo, Inches(5.16), Inches(5.05), Inches(3.0), Inches(1.65))
    print('Slide  1: Title')


def slide_outline(prs):
    slide = white_slide(prs)
    slide_header(slide, 'Presentation Outline')
    sections = [
        ('1.  The Problem',
         ['Why FPGA for tracking?', 'Frequency-domain insight', 'NCC / MOSSE / KCF comparison']),
        ('2.  KCF Algorithm',
         ['Hann window & Gaussian label', 'Training: alpha = Y/(|X|^2 + lambda)', 'Detection & response map']),
        ('3.  Live Demo',
         ['Frame 1 — target selection + training', 'Frame 2 — NCC full-image search', 'Frame 3 — KCF-only tracking']),
        ('4.  Hardware & Integration',
         ['7-stage RTL pipeline (~19,300 cycles)', 'AXI interface & processor integration', 'Future work']),
    ]
    xs = [Inches(0.3), Inches(6.85)]
    ys = [Inches(1.25), Inches(4.1)]
    for i, (sec_title, items) in enumerate(sections):
        card(slide, xs[i % 2], ys[i // 2], Inches(6.18), Inches(2.6),
             sec_title, items, title_size=16, body_size=14)
    print('Slide  2: Outline')


def slide_problem(prs):
    slide = white_slide(prs)
    slide_header(slide, 'Why FPGA? The Frequency-Domain Insight')
    bullets(slide, [
        ('Challenge:',
         'Real-time tracking on edge devices needs low latency and low power '
         '— CPUs execute serially and are too slow for high-frame-rate video.'),
        ('Key Insight:',
         'KCF exploits that all cyclic shifts of a patch form a circulant matrix '
         'diagonalised by the DFT — tracking reduces to element-wise division: '
         'alpha = Y / (|X|^2 + lambda)'),
        ('FPGA Advantage:',
         'Dedicated parallel dataflow maps directly to the FFT butterfly pipeline. '
         'Fixed-point arithmetic eliminates FPUs, cutting LUT usage significantly.'),
    ], x=Inches(0.35), y=Inches(1.3), w=Inches(5.7), h=Inches(5.6), size=15, leading=14)
    pic(slide, 'KCF.png',       Inches(6.2), Inches(1.2),  Inches(6.8), Inches(3.3))
    pic(slide, 'fig_response_map.png', Inches(6.2), Inches(4.6), Inches(6.8), Inches(2.65))
    print('Slide  4: Problem & Insight')


def slide_algo_table(prs):
    slide = white_slide(prs)
    slide_header(slide, 'Algorithm Survey: NCC vs MOSSE vs KCF')
    headers = ['Property', 'NCC', 'MOSSE', 'KCF']
    rows = [
        ['Purpose',       'Target acquisition',  'Fast tracking',     'Robust per-frame tracking'],
        ['Feature space', 'Raw pixels',           'Raw pixels',        'Kernel (linear / RBF)'],
        ['Complexity',    'O(MN log MN)',         'O(N^2 log N)',      'O(N^2 log N)'],
        ['Robustness',    'Low',                  'Medium',            'High'],
        ['Update rule',   'None',                 'Online average',    'Learning rate eta'],
        ['HW cost',       'Medium (1 FFT pair)',  'Low (2 FFT)',       'Medium (6 FFT)'],
        ['This project',  'Frame 2 only',         'Algorithm basis',   'Frames 3+'],
    ]
    tbl(slide, headers, rows, x=Inches(0.35), y=Inches(1.25),
        w=Inches(12.6), h=Inches(5.85), font_size=15)
    print('Slide  5: Algorithm table')


def slide_training(prs):
    slide = white_slide(prs)
    slide_header(slide, 'KCF Training: Hann Window & Gaussian Label')
    pic(slide, 'fig_pres_hann_effect.png',    Inches(0.3), Inches(1.2),  Inches(5.9), Inches(2.9))
    pic(slide, 'fig_pres_gaussian_label.png', Inches(0.3), Inches(4.2),  Inches(5.9), Inches(3.05))
    eq = slide.shapes.add_shape(1, Inches(6.5), Inches(1.2), Inches(6.5), Inches(1.0))
    _fill(eq, LGRAY); _border(eq, NAVY, 1.5)
    tb(slide, 'alpha = Y  /  ( |X|^2  +  lambda )',
       Inches(6.6), Inches(1.25), Inches(6.3), Inches(0.85),
       size=24, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    bullets(slide, [
        ('1.', 'Extract 32x32 patch x from Frame 1'),
        ('2.', 'Apply Hann window:  x_w = x * w^2'),
        ('3.', 'Compute FFT:  X = FFT2(x_w)'),
        ('4.', 'Gaussian label:  Y = FFT2(y,  sigma=2.0)'),
        ('5.', 'alpha = Y / (|X|^2 + lambda),  lambda=0.01'),
        ('6.', 'Store alpha in ROM — computed once offline'),
    ], x=Inches(6.5), y=Inches(2.38), w=Inches(6.5), h=Inches(4.8), size=15, leading=10, bullet='>')
    print('Slide  7: Training')


def slide_detection(prs):
    slide = white_slide(prs)
    slide_header(slide, 'KCF Detection: Response Map & Peak Finding')
    eq = slide.shapes.add_shape(1, Inches(0.35), Inches(1.2), Inches(6.15), Inches(1.0))
    _fill(eq, LGRAY); _border(eq, NAVY, 1.5)
    tb(slide, 'R  =  IFFT2( conj(X) . Z . alpha )',
       Inches(0.45), Inches(1.25), Inches(5.95), Inches(0.85),
       size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    bullets(slide, [
        ('1.', 'Extract 32x32 patch z around last known position'),
        ('2.', 'Apply Hann window:  z_w = z * w^2'),
        ('3.', 'Compute FFT:  Z = FFT2(z_w)'),
        ('4.', 'Correlate:  R_hat = conj(X) . Z . alpha  (element-wise)'),
        ('5.', 'IFFT via conjugate trick:  R = IFFT2(R_hat)'),
        ('6.', 'Peak of R gives displacement (d_row, d_col) -> new target position'),
    ], x=Inches(0.35), y=Inches(2.38), w=Inches(6.15), h=Inches(4.8), size=15, leading=10, bullet='>')
    pic(slide, 'fig_response_map.png',    Inches(6.65), Inches(1.2),  Inches(6.4), Inches(3.2))
    pic(slide, 'fig_pres_response_3d.png',Inches(6.65), Inches(4.5),  Inches(6.4), Inches(2.75))
    print('Slide  8: Detection')


def slide_demo_frame1(prs):
    slide = white_slide(prs)
    slide_header(slide, 'Frame 1 — Target Selection & KCF Training', subtitle='Step 1-2 / 7')
    pic(slide, 'demo_s1_select.png',  Inches(0.25), Inches(1.2), Inches(6.45), Inches(5.85))
    pic(slide, 'demo_s2_trained.png', Inches(6.85), Inches(1.2), Inches(6.2),  Inches(5.85))
    cap = slide.shapes.add_shape(1, Inches(0.25), Inches(7.1), Inches(12.8), Inches(0.26))
    _fill(cap, LGRAY); _border(cap, SILVER, 0.3)
    tb(slide, 'User clicks to select a 32x32 target patch  ->  KCF filter alpha computed offline and stored in ROM',
       Inches(0.4), Inches(7.11), Inches(12.5), Inches(0.26),
       size=12, color=col('444444'), align=PP_ALIGN.CENTER)
    print('Slide 10: Demo Frame 1')


def slide_demo_frame2(prs):
    slide = white_slide(prs)
    slide_header(slide, 'Frame 2 — NCC Full-Image Search', subtitle='Step 4 / 7')
    pic(slide, 'demo_s3_ncc.png', Inches(0.25), Inches(1.2), Inches(12.8), Inches(5.55))
    cap = slide.shapes.add_shape(1, Inches(0.25), Inches(6.85), Inches(12.8), Inches(0.42))
    _fill(cap, LGRAY); _border(cap, SILVER, 0.3)
    tb(slide,
       'NCC score = 1.000  |  Target found at (146, 853)  |  '
       'NCC used ONCE for initial acquisition — too slow for every frame',
       Inches(0.4), Inches(6.88), Inches(12.5), Inches(0.4),
       size=13, color=col('333333'), align=PP_ALIGN.CENTER)
    print('Slide 11: Demo Frame 2')


def slide_demo_frame3(prs):
    slide = white_slide(prs)
    slide_header(slide, 'NCC -> KCF Handoff & Frame 3 Tracking', subtitle='Steps 5-6 / 7')
    pic(slide, 'demo_s4_handoff.png', Inches(0.25), Inches(1.2),  Inches(12.8), Inches(2.55))
    pic(slide, 'demo_s5_kcf.png',     Inches(0.25), Inches(3.88), Inches(12.8), Inches(3.38))
    print('Slide 12: Demo Frame 3 & Handoff')


def slide_results(prs):
    slide = white_slide(prs)
    slide_header(slide, 'Results Summary', subtitle='Response maps + hardware simulation')
    pic(slide, 'demo_s6_summary.png', Inches(0.25), Inches(1.2), Inches(7.85), Inches(5.85))
    pic(slide, 'waveform.png',         Inches(8.3),  Inches(1.2), Inches(4.75), Inches(2.75))
    stats = slide.shapes.add_shape(1, Inches(8.3), Inches(4.1), Inches(4.75), Inches(2.9))
    _fill(stats, LGRAY); _border(stats, NAVY, 1.2)
    tb(slide, 'Hardware Simulation',
       Inches(8.42), Inches(4.18), Inches(4.5), Inches(0.38),
       size=14, bold=True, color=NAVY)
    bullets(slide, [
        'done pulse confirmed in Vivado',
        '~19,300 cycles = 193 us @ 100 MHz',
        'peak_row / peak_col match Python ref',
        'Hardware-software agreement verified',
    ], x=Inches(8.42), y=Inches(4.58), w=Inches(4.5), h=Inches(2.3),
       size=13, color=DTEXT, leading=5, bullet='v')
    print('Slide 13: Results')


def slide_hardware(prs):
    slide = white_slide(prs)
    slide_header(slide, 'RTL Pipeline & Processor Integration')
    tb(slide, '7-Stage Detection Pipeline',
       Inches(0.35), Inches(1.2), Inches(6.1), Inches(0.42),
       size=17, bold=True, color=NAVY)
    tbl(slide,
        ['Stage', 'Cycles', 'Operation'],
        [['S_HANN_WIN',  '1,024',   'Hann window multiply'],
         ['S_FFT_RUN',   '~5,600',  '2D-FFT (32 row + 32 col)'],
         ['S_ELEM_MUL',  '1,024',   'conj(X).Z.alpha element-wise'],
         ['S_IFFT_RUN',  '~5,600',  '2D-IFFT (conjugate trick)'],
         ['S_PEAK_RUN',  '1,024',   'Scan for response peak'],
         ['TOTAL',       '~19,300', '~193 us @ 100 MHz']],
        x=Inches(0.35), y=Inches(1.68), w=Inches(6.1), h=Inches(4.6), font_size=13)

    tb(slide, 'AXI Interface — Processor Integration',
       Inches(6.75), Inches(1.2), Inches(6.2), Inches(0.42),
       size=17, bold=True, color=NAVY)
    tbl(slide,
        ['Offset', 'Register', 'Access'],
        [['0x00', 'CTRL',         'W — trigger start'],
         ['0x04', 'THRESHOLD',    'R/W — confidence level'],
         ['0x08', 'PEAK_ROW',     'R — row index (0-31)'],
         ['0x0C', 'PEAK_COL',     'R — col index (0-31)'],
         ['0x10', 'PEAK_VAL',     'R — response magnitude'],
         ['0x14', 'TARGET_FOUND', 'R — valid detection bit']],
        x=Inches(6.75), y=Inches(1.68), w=Inches(6.2), h=Inches(3.38), font_size=12)

    ib = slide.shapes.add_shape(1, Inches(6.75), Inches(5.15), Inches(6.2), Inches(1.12))
    _fill(ib, LGRAY); _border(ib, NAVY, 1.0)
    tb(slide,
       'AXI4-Stream: 1024 x 16-bit pixel stream  (tready / tvalid / tlast)\n'
       'AXI4-Lite: control + result registers\n'
       'Integration with processor via AXI interconnect',
       Inches(6.87), Inches(5.2), Inches(5.96), Inches(1.05), size=13, color=DTEXT)
    print('Slide 15: Hardware & Integration')


def slide_future(prs):
    slide = white_slide(prs)
    slide_header(slide, 'Future Work & Thank You')
    tb(slide, 'Future Work',
       Inches(0.35), Inches(1.2), Inches(6.0), Inches(0.42),
       size=20, bold=True, color=NAVY)
    items = [
        ('FPGA Synthesis',    'Artix-7 LUT/FF/BRAM utilisation + timing closure at 100 MHz'),
        ('64x64 Resolution',  '6-stage FFT for higher accuracy on larger targets'),
        ('Gaussian Kernel',   'Hardware exponential LUT for RBF kernel — improves robustness'),
        ('Online Update',     'Adaptive filter: alpha_new = (1-eta)*alpha_old + eta*alpha_frame'),
        ('Camera Pipeline',   'MIPI CSI-2 / OV7670 input + HDMI output on Nexys A7'),
    ]
    for i, (head, body) in enumerate(items):
        y = Inches(1.78) + i * Inches(1.0)
        fc = slide.shapes.add_shape(1, Inches(0.35), y, Inches(6.05), Inches(0.88))
        _fill(fc, LGRAY); _border(fc, BLUE, 0.7)
        tb(slide, f'>  {head}', Inches(0.5), y + Inches(0.04), Inches(5.75), Inches(0.34),
           size=14, bold=True, color=NAVY)
        tb(slide, body, Inches(0.5), y + Inches(0.42), Inches(5.75), Inches(0.42),
           size=12, color=DTEXT)

    ty = slide.shapes.add_shape(1, Inches(6.72), Inches(1.2), Inches(6.28), Inches(5.2))
    _fill(ty, NAVY); ty.line.fill.background()
    tb(slide, 'Thank You',
       Inches(6.82), Inches(2.1), Inches(6.08), Inches(1.0),
       size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(slide, 'Questions?',
       Inches(6.82), Inches(3.08), Inches(6.08), Inches(0.65),
       size=26, color=col('A0BAD8'), align=PP_ALIGN.CENTER)
    hbar(slide, Inches(3.83), x=Inches(7.3), w=Inches(5.1), color=BLUE, height=Inches(0.03))
    tb(slide, 'Akhil Sriram  |  2210110134',
       Inches(6.82), Inches(3.95), Inches(6.08), Inches(0.4),
       size=14, color=WHITE, align=PP_ALIGN.CENTER)
    tb(slide, 'Dr. Venkatnarayan Hariharan',
       Inches(6.82), Inches(4.35), Inches(6.08), Inches(0.38),
       size=13, color=col('A0BAD8'), align=PP_ALIGN.CENTER)
    tb(slide, 'Shiv Nadar Institution of Eminence',
       Inches(6.82), Inches(4.73), Inches(6.08), Inches(0.35),
       size=12, color=col('8090B0'), align=PP_ALIGN.CENTER)
    print('Slide 16: Future Work & Thank You')


# ── Main ───────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    slide_title(prs)                                    #  1
    slide_outline(prs)                                  #  2
    sec(prs, 'The Visual Tracking Problem')             #  3
    slide_problem(prs)                                  #  4
    slide_algo_table(prs)                               #  5
    sec(prs, 'KCF Algorithm')                           #  6
    slide_training(prs)                                 #  7
    slide_detection(prs)                                #  8
    sec(prs, 'Live Demo: NCC to KCF Tracking Pipeline') #  9
    slide_demo_frame1(prs)                              # 10
    slide_demo_frame2(prs)                              # 11
    slide_demo_frame3(prs)                              # 12
    slide_results(prs)                                  # 13
    sec(prs, 'RTL Hardware & Integration')              # 14
    slide_hardware(prs)                                 # 15
    slide_future(prs)                                   # 16

    prs.save(OUT_PATH)
    print(f'\nSaved: {OUT_PATH}  ({len(prs.slides)} slides)')


if __name__ == '__main__':
    main()
